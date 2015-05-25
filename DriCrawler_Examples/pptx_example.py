from pptx import Presentation

path_to_presentation = "/Users/morrj140/Development/GitRepository/DirCrawler/Business Resource Estimates.pptx"
prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_textframe:
            continue
        for paragraph in shape.textframe.paragraphs:
            for run in paragraph.runs:
                print run.text
                text_runs.append(run.text)
