import unicodedata

from nl_lib import Logger
from nl_lib.Concepts import Concepts
import openxmllib
import xlrd
from pptx import Presentation

from examples import docx


logger = Logger.setupLogging(__name__)

def getPPTXText(filename, ftype):
    logger.info("filename: %s" % filename)

    prs = Presentation(filename)

    newparatextlist = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_textframe:
                continue
            for paragraph in shape.textframe.paragraphs:
                for run in paragraph.runs:
                    logger.info("PPTX : %s" % run.text)
                    if run.text != None:
                        #`unicodedata.normalize('NFKD', run.text).encode('ascii', 'ignore')
                        newparatextlist.append(run.text + ". ")
    return newparatextlist

def getXLSText(filename, ftype):
    logger.debug("filename: %s" % filename)

    newparatextlist = []
    workbook = xlrd.open_workbook(filename)

    #sheet = "Specific Requirements"
    #worksheet = workbook.sheet_by_name(sheet)

    CellTypes = ["Empty", "Text", "Number", "Date", "Boolean", "Error", "Blank"]

    for worksheet_name in workbook.sheet_names():
        worksheet = workbook.sheet_by_name(worksheet_name)
        num_rows = worksheet.nrows - 1
        num_cells = worksheet.ncols - 1
        curr_row = -1

        while curr_row < num_rows:
            curr_row += 1
            row = worksheet.row(curr_row)
            logger.info('Row: %s' % curr_row)
            curr_cell = -1
            while curr_cell < num_cells:
                curr_cell += 1
                # Cell Types: 0=Empty, 1=Text, 2=Number, 3=Date, 4=Boolean, 5=Error, 6=Blank
                cell_type = worksheet.cell_type(curr_row, curr_cell)
                cell_value = worksheet.cell_value(curr_row, curr_cell)
                if cell_type == 1:
                    unicodedata.normalize('NFKD', cell_value).encode('ascii', 'ignore')
                    logger.info("XLXS : %s" % cell_value)
                    newparatextlist.append(cell_value + ". ")

    return newparatextlist

def getDOCXText(filename, ftype):
        logger.info("filename: %s" % filename)

        document = docx.opendocx(filename)

        # Fetch all the text out of the document we just created
        paratextlist = docx.getdocumenttext(document)

        # Make explicit unicode version
        newparatextlist = []
        for paratext in paratextlist:
            unicodedata.normalize('NFKD', paratext).encode('ascii', 'ignore')
            logger.info("DOCX : %s" % paratext)
            newparatextlist.append(paratext + ". ")

        return newparatextlist

def getOpenXmlText(filename, ftype):
        logger.info("OpenXmlText: %s" % filename)

        document = openxmldoc

        doc = openxmllib.openXmlDocument(path=filename)
        c = Concepts(filename, ftype)

        logger.debug ("%s\n" % (doc.allProperties))

        ap = c.addConceptKeyType("allProperties","PROPERTIES")
        for x in doc.allProperties:
            logger.info("cp %s:%s" % (x, doc.allProperties[x]))
            ap.addConceptKeyType(doc.allProperties[x], x)

        logger.info("it %s\n" % (doc.indexableText(include_properties=True)))
        c.addConceptKeyType(doc.indexableText(include_properties=True),"TEXT")

        return c

if __name__ == "__main__":
    #concepts = getOpenXmlText(“example.docx", "DOCX")

    #l = getDOCXText("example.docx", "DOCX")

    #l = getPPTXText("example.pptx", "PPTX")

    l =getXLSText("example.xlsx", "XLSX")

    #logger.info ("log of concepts")
    #concepts.logConcepts()